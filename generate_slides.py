import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add Background Image
    img_path = r'C:\Users\User\.gemini\antigravity\brain\90ad7ee2-4f26-4001-833b-dca290e80550\presentation_title_slide_background_1773608932043.png'
    if os.path.exists(img_path):
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # --- Slide 2 - 12: Content ---
    content = [
        ("The Problem Statement", [
            "Manual pentesting is slow and expensive.",
            "Automated tools produce mountains of raw data.",
            "Beginners struggle to interpret 'low-level' results.",
            "Turning data into actionable security insights."
        ]),
        ("Our Objectives", [
            "Automate reconnaissance and scanning.",
            "Apply AI intelligence to explain findings.",
            "Create an accessible and professional UI.",
            "Generate actionable, boardroom-ready reports."
        ]),
        ("Key Features", [
            "Dynamic Recon: Port, DNS, SSL, and Tech scanning.",
            "AI-Powered Interpretation: Llama 3.2 analysis.",
            "Interactive Assistant: Built-in security tutor.",
            "Professional Reporting: Instant Markdown/PDF."
        ]),
        ("Technical Architecture", [
            "Backend: FastAPI (Python).",
            "AI Core: Ollama (Llama 3.2) - Private & Local.",
            "Scanning Engine: Custom Python Modules.",
            "Frontend: Modern Vanilla JS + Glassmorphic CSS."
        ]),
        ("Deep Visibility (Recon)", [
            "Port Scanner: Service discovery.",
            "DNS Enumeration: Infrastructure visibility.",
            "Security Header Audit: HSTS, CSP checks.",
            "SSL/TLS Status: Encryption health."
        ]),
        ("The Intelligence Layer", [
            "Uses scan data as context for AI interpretation.",
            "Identifies and explains chainable vulnerabilities.",
            "Ranks findings by severity (Critical to Info).",
            "Provides educational remediation steps."
        ]),
        ("Professional UI/UX", [
            "Premium Design: Glassmorphism and dark mode.",
            "User-Centric: Dashboard overview.",
            "Fully Responsive: Works on all devices.",
            "Interactive Chat: Floating assistant modal."
        ]),
        ("Actionable Reporting", [
            "One-click professional report generation.",
            "Markdown export for developer collaboration.",
            "Print-ready PDF styling for executive review.",
            "AI-built remediation priority matrix."
        ]),
        ("Value Proposition", [
            "Students: Educational bridge to professional tools.",
            "Startups: Faster, low-cost surface analysis.",
            "Educators: Platform for security demonstrations.",
            "Efficiency: Reduces Recon-to-Report time by 80%."
        ]),
        ("The Roadmap", [
            "Integration with Nuclei templates (10k+ checks).",
            "Cloud-based scalable infrastructure.",
            "Multi-target bulk scanning and scheduling.",
            "Deep-scan AI focused on OWASP Top 10."
        ]),
        ("Conclusion & Q&A", [
            "AI as a force multiplier for security analysts.",
            "Smart, Fast, and Accessible security auditing.",
            "Contact: [Your Contact Info]",
            "Thank you for your time!"
        ])
    ]

    for title_text, bullets in content:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title = slide.shapes.title
        title.text = title_text
        
        # Set Content
        tf = slide.placeholders[1].text_frame
        tf.text = bullets[0]
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    prs.save('Hackathon_Presentation.pptx')
    print("Presentation created successfully: Hackathon_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
